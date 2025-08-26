"""
Enhanced Document Processor for Multi-Format Support

Supports: .pdf, .docx, .txt, .ppt, .pptx, .xlsx, .csv, .md, and SQL databases
"""

from __future__ import annotations
import os
import sys
import pandas as pd
import sqlite3
from pathlib import Path
from typing import Iterable, List, Dict, Any, Optional, Union
from fastapi import UploadFile
from langchain.schema import Document
from langchain_community.document_loaders import (
    PyPDFLoader, 
    Docx2txtLoader, 
    TextLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    CSVLoader,
    UnstructuredMarkdownLoader
)
from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException

# Enhanced supported extensions
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".ppt", ".pptx", 
    ".xlsx", ".csv", ".md", ".sql", ".db"
}

class EnhancedDocumentProcessor:
    """
    Enhanced document processor supporting multiple file formats
    including tables and image data extraction
    """
    
    def __init__(self):
        self.logger = log
    
    def load_documents(self, paths: Iterable[Path]) -> List[Document]:
        """Load documents using appropriate loader based on extension."""
        docs: List[Document] = []
        
        try:
            for file_path in paths:
                ext = file_path.suffix.lower()
                
                if ext not in SUPPORTED_EXTENSIONS:
                    self.logger.warning(f"Unsupported extension skipped: {file_path}")
                    continue
                
                try:
                    if ext == ".pdf":
                        loader = PyPDFLoader(str(file_path))
                    elif ext == ".docx":
                        loader = Docx2txtLoader(str(file_path))
                    elif ext == ".txt":
                        loader = TextLoader(str(file_path), encoding="utf-8")
                    elif ext in [".ppt", ".pptx"]:
                        try:
                            loader = UnstructuredPowerPointLoader(str(file_path))
                        except Exception as e:
                            self.logger.warning(f"UnstructuredPowerPointLoader failed, trying alternative: {str(e)}")
                            # Fallback to python-pptx
                            from pptx import Presentation
                            prs = Presentation(str(file_path))
                            text_content = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text_content.append(shape.text)
                            content = "\n".join(text_content)
                            doc = Document(
                                page_content=content,
                                metadata={
                                    "source": str(file_path),
                                    "file_type": "powerpoint",
                                    "slides": len(prs.slides)
                                }
                            )
                            docs.append(doc)
                            continue
                    elif ext == ".xlsx":
                        loader = UnstructuredExcelLoader(str(file_path))
                    elif ext == ".csv":
                        loader = CSVLoader(str(file_path))
                    elif ext == ".md":
                        loader = UnstructuredMarkdownLoader(str(file_path))
                    elif ext in [".sql", ".db"]:
                        # Handle SQL files and databases
                        docs.extend(self._process_sql_file(file_path))
                        continue
                    else:
                        self.logger.warning(f"Extension {ext} not yet implemented")
                        continue
                    
                    docs.extend(loader.load())
                    self.logger.info(f"Successfully loaded: {file_path}")
                    
                except Exception as e:
                    self.logger.error(f"Failed to load {file_path}: {str(e)}")
                    continue
            
            self.logger.info(f"Total documents loaded: {len(docs)}")
            return docs
            
        except Exception as e:
            self.logger.error(f"Failed loading documents: {str(e)}")
            raise DocumentPortalException("Error loading documents", e) from e
    
    def _process_sql_file(self, file_path: Path) -> List[Document]:
        """Process SQL files and databases"""
        docs = []
        
        try:
            if file_path.suffix.lower() == ".sql":
                # Handle SQL script files
                with open(file_path, 'r', encoding='utf-8') as f:
                    sql_content = f.read()
                
                doc = Document(
                    page_content=sql_content,
                    metadata={
                        "source": str(file_path),
                        "file_type": "sql_script",
                        "tables": self._extract_table_names_from_sql(sql_content)
                    }
                )
                docs.append(doc)
                
            elif file_path.suffix.lower() == ".db":
                # Handle SQLite databases
                docs.extend(self._process_sqlite_database(file_path))
                
        except Exception as e:
            self.logger.error(f"Failed to process SQL file {file_path}: {str(e)}")
        
        return docs
    
    def _process_sqlite_database(self, db_path: Path) -> List[Document]:
        """Process SQLite database and extract table schemas and data"""
        docs = []
        
        try:
            conn = sqlite3.connect(db_path)
            cursor = conn.cursor()
            
            # Get all table names
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            for table in tables:
                table_name = table[0]
                
                # Get table schema
                cursor.execute(f"PRAGMA table_info({table_name});")
                schema = cursor.fetchall()
                
                # Get sample data (first 10 rows)
                cursor.execute(f"SELECT * FROM {table_name} LIMIT 10;")
                sample_data = cursor.fetchall()
                
                # Get column names
                cursor.execute(f"SELECT * FROM {table_name} LIMIT 0;")
                columns = [description[0] for description in cursor.description]
                
                # Create document content
                content = f"Table: {table_name}\n\n"
                content += "Schema:\n"
                for col in schema:
                    content += f"  {col[1]} ({col[2]})\n"
                
                if sample_data:
                    content += "\nSample Data:\n"
                    content += " | ".join(columns) + "\n"
                    content += "-" * (len(" | ".join(columns))) + "\n"
                    for row in sample_data:
                        content += " | ".join(str(cell) for cell in row) + "\n"
                
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": str(db_path),
                        "file_type": "sqlite_database",
                        "table_name": table_name,
                        "columns": columns,
                        "row_count": len(sample_data)
                    }
                )
                docs.append(doc)
            
            conn.close()
            
        except Exception as e:
            self.logger.error(f"Failed to process SQLite database {db_path}: {str(e)}")
        
        return docs
    
    def _extract_table_names_from_sql(self, sql_content: str) -> List[str]:
        """Extract table names from SQL script"""
        import re
        # Simple regex to find table names in CREATE TABLE statements
        table_pattern = r'CREATE\s+TABLE\s+(?:IF\s+NOT\s+EXISTS\s+)?([a-zA-Z_][a-zA-Z0-9_]*)'
        tables = re.findall(table_pattern, sql_content, re.IGNORECASE)
        return tables
    
    def extract_tables_from_documents(self, docs: List[Document]) -> Dict[str, pd.DataFrame]:
        """Extract table data from documents"""
        tables = {}
        
        for doc in docs:
            source = doc.metadata.get("source", "unknown")
            content = doc.page_content
            
            # Look for CSV-like content in the document
            lines = content.split('\n')
            table_data = []
            
            for line in lines:
                if ',' in line and len(line.split(',')) > 2:
                    table_data.append(line.split(','))
            
            if table_data:
                try:
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    tables[f"table_{source}"] = df
                except Exception as e:
                    self.logger.warning(f"Failed to create DataFrame from {source}: {str(e)}")
        
        return tables
    
    def extract_images_from_documents(self, docs: List[Document]) -> Dict[str, List[str]]:
        """Extract image references from documents"""
        images = {}
        
        for doc in docs:
            source = doc.metadata.get("source", "unknown")
            content = doc.page_content
            
            # Look for image references (basic pattern matching)
            import re
            image_patterns = [
                r'!\[.*?\]\((.*?)\)',  # Markdown images
                r'<img[^>]+src=["\']([^"\']+)["\']',  # HTML images
                r'\.(jpg|jpeg|png|gif|bmp|svg)',  # Image file extensions
            ]
            
            found_images = []
            for pattern in image_patterns:
                matches = re.findall(pattern, content, re.IGNORECASE)
                found_images.extend(matches)
            
            if found_images:
                images[source] = found_images
        
        return images

class FastAPIFileAdapter:
    """Enhanced adapter for FastAPI UploadFile with multi-format support"""
    
    SUPPORTED_EXTENSIONS = SUPPORTED_EXTENSIONS
    
    def __init__(self, uf: UploadFile):
        self._uf = uf
        self.name = uf.filename
        self.content_type = uf.content_type
    
    def getbuffer(self) -> bytes:
        self._uf.file.seek(0)
        return self._uf.file.read()
    
    def get_extension(self) -> str:
        """Get file extension"""
        return Path(self.name).suffix.lower() if self.name else ""
    
    def is_supported(self) -> bool:
        """Check if file type is supported"""
        return self.get_extension() in SUPPORTED_EXTENSIONS

def concat_for_analysis(docs: List[Document]) -> str:
    """Concatenate documents for analysis with enhanced metadata"""
    parts = []
    for doc in docs:
        src = doc.metadata.get("source") or doc.metadata.get("file_path") or "unknown"
        file_type = doc.metadata.get("file_type", "unknown")
        parts.append(f"\n--- SOURCE: {src} (Type: {file_type}) ---\n{doc.page_content}")
    return "\n".join(parts)

def concat_for_comparison(ref_docs: List[Document], act_docs: List[Document]) -> str:
    """Concatenate documents for comparison"""
    left = concat_for_analysis(ref_docs)
    right = concat_for_analysis(act_docs)
    return f"<<REFERENCE_DOCUMENTS>>\n{left}\n\n<<ACTUAL_DOCUMENTS>>\n{right}"
